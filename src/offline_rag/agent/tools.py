from __future__ import annotations

import asyncio
import shutil
import sys
import tempfile
from pathlib import Path

import httpx

# ---------------------------------------------------------------------------
# Tool definitions (sent to Ollama for function calling)
# ---------------------------------------------------------------------------

TOOL_DEFS: list[dict] = [
    {
        "type": "function",
        "function": {
            "name": "web_search",
            "description": (
                "Search the internet using DuckDuckGo. Use this to find information, "
                "tutorials, documentation, or how to do something you don't know."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Search query"},
                    "max_results": {
                        "type": "integer",
                        "description": "Number of results to return (default 6)",
                        "default": 6,
                    },
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "web_fetch",
            "description": "Fetch and read the full text content of a webpage.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "URL to fetch"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_python",
            "description": (
                "Execute Python code on the user's machine. "
                "The user will be asked to confirm before it runs."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "code": {"type": "string", "description": "Python code to execute"},
                    "description": {
                        "type": "string",
                        "description": "One-line plain-English description of what this code does",
                    },
                },
                "required": ["code", "description"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_shell",
            "description": (
                "Execute a PowerShell command on the user's Windows machine. "
                "The user will be asked to confirm before it runs."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {"type": "string", "description": "PowerShell command"},
                    "description": {
                        "type": "string",
                        "description": "One-line plain-English description of what this command does",
                    },
                },
                "required": ["command", "description"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "Read the text contents of a file.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Absolute or ~ path to the file",
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": (
                "Write (create or overwrite) a file with the given content. "
                "The user will be asked to confirm before writing."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Destination file path"},
                    "content": {"type": "string", "description": "Full file content to write"},
                    "description": {
                        "type": "string",
                        "description": "One-line description of what this file contains",
                    },
                },
                "required": ["path", "content", "description"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_files",
            "description": "List files in a directory, with optional glob filtering.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Directory path"},
                    "pattern": {
                        "type": "string",
                        "description": "Glob pattern, e.g. *.py (default: *)",
                        "default": "*",
                    },
                    "recursive": {
                        "type": "boolean",
                        "description": "Search recursively (default false)",
                        "default": False,
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "move_file",
            "description": (
                "Move or rename a file or directory. "
                "The user will be asked to confirm before moving."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "src": {"type": "string", "description": "Source path"},
                    "dst": {"type": "string", "description": "Destination path"},
                },
                "required": ["src", "dst"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "delete_file",
            "description": (
                "Permanently delete a file or directory. Cannot be undone. "
                "The user will be asked to confirm before deleting."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to delete"},
                    "recursive": {
                        "type": "boolean",
                        "description": "Delete directory and all its contents (default false)",
                        "default": False,
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "make_slides",
            "description": "Create a presentation (.pptx PowerPoint or .pdf) with the given slides.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Presentation title"},
                    "slides": {
                        "type": "array",
                        "description": "Ordered list of slides",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string", "description": "Slide heading"},
                                "content": {
                                    "type": "string",
                                    "description": "Slide body — one bullet point per line",
                                },
                            },
                            "required": ["title", "content"],
                        },
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output file path, e.g. ~/Desktop/slides.pptx",
                    },
                    "format": {
                        "type": "string",
                        "enum": ["pptx", "pdf"],
                        "description": "File format (default: pptx)",
                        "default": "pptx",
                    },
                },
                "required": ["title", "slides", "output_path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "rag_search",
            "description": "Search the user's locally indexed personal documents.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Search query"},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "install_package",
            "description": (
                "Install a Python package via pip. "
                "The user will be asked to confirm before installing."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "package": {
                        "type": "string",
                        "description": "Package name, e.g. 'requests' or 'pandas==2.0'",
                    },
                },
                "required": ["package"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "index_documents",
            "description": "Index new documents so they become searchable via rag_search.",
            "parameters": {
                "type": "object",
                "properties": {
                    "paths": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "File or directory paths to index",
                    },
                },
                "required": ["paths"],
            },
        },
    },
]

# Tools that require user confirmation before executing
RISKY_TOOLS: frozenset[str] = frozenset(
    {"run_python", "run_shell", "write_file", "move_file", "delete_file", "install_package"}
)


# ---------------------------------------------------------------------------
# Implementations
# ---------------------------------------------------------------------------

async def web_search(query: str, max_results: int = 6) -> str:
    try:
        from duckduckgo_search import DDGS
    except ImportError:
        return "ERROR: duckduckgo-search is not installed. Run: pip install duckduckgo-search"

    def _sync_search() -> list[dict]:
        with DDGS() as ddgs:
            return list(ddgs.text(query, max_results=max_results))

    try:
        results = await asyncio.get_event_loop().run_in_executor(None, _sync_search)
    except Exception as e:
        return f"Search failed: {e}"

    if not results:
        return "No results found."

    lines: list[str] = []
    for r in results:
        lines.append(f"### {r.get('title', 'No title')}")
        lines.append(r.get("href", ""))
        lines.append(r.get("body", ""))
        lines.append("")
    return "\n".join(lines)


async def web_fetch(url: str) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return "ERROR: beautifulsoup4 is not installed. Run: pip install beautifulsoup4"

    try:
        async with httpx.AsyncClient(
            timeout=30,
            follow_redirects=True,
            headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"},
        ) as client:
            resp = await client.get(url)
            resp.raise_for_status()
            html = resp.text
    except Exception as e:
        return f"Failed to fetch {url}: {e}"

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header", "aside", "noscript"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    lines = [ln for ln in text.split("\n") if ln.strip()]
    text = "\n".join(lines)

    if len(text) > 12000:
        text = text[:12000] + "\n...[truncated — page is long]"

    return text or "(page appears empty or is JavaScript-rendered)"


async def run_python(code: str, description: str) -> str:
    with tempfile.NamedTemporaryFile(
        mode="w", suffix=".py", delete=False, encoding="utf-8"
    ) as f:
        f.write(code)
        tmp = Path(f.name)

    try:
        proc = await asyncio.create_subprocess_exec(
            sys.executable,
            str(tmp),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=str(Path.home()),
        )
        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=120)
        except asyncio.TimeoutError:
            proc.kill()
            return "Timeout: Python script ran for more than 120 seconds and was killed."

        output = stdout.decode(errors="replace")
        err = stderr.decode(errors="replace")
        if err:
            output += f"\nSTDERR:\n{err}"
        if proc.returncode != 0:
            output += f"\nExit code: {proc.returncode}"
        return output or "(no output)"
    finally:
        tmp.unlink(missing_ok=True)


async def run_shell(command: str, description: str) -> str:
    proc = await asyncio.create_subprocess_exec(
        "powershell",
        "-NoProfile",
        "-NonInteractive",
        "-Command",
        command,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
        cwd=str(Path.home()),
    )
    try:
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=120)
    except asyncio.TimeoutError:
        proc.kill()
        return "Timeout: command ran for more than 120 seconds and was killed."

    output = stdout.decode(errors="replace")
    err = stderr.decode(errors="replace")
    if err:
        output += f"\nSTDERR:\n{err}"
    if proc.returncode != 0:
        output += f"\nExit code: {proc.returncode}"
    return output or "(no output)"


def read_file(path: str) -> str:
    p = Path(path).expanduser()
    if not p.exists():
        return f"File not found: {p}"
    if not p.is_file():
        return f"Not a file: {p}"
    try:
        text = p.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"Could not read file: {e}"
    if len(text) > 20000:
        text = text[:20000] + "\n...[truncated]"
    return text


def write_file(path: str, content: str, description: str) -> str:
    p = Path(path).expanduser()
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content, encoding="utf-8")
    return f"Written {len(content):,} chars ({p.stat().st_size:,} bytes) to {p}"


def list_files(path: str, pattern: str = "*", recursive: bool = False) -> str:
    p = Path(path).expanduser()
    if not p.exists():
        return f"Path not found: {p}"
    if not p.is_dir():
        return f"Not a directory: {p}"

    matches = sorted(p.rglob(pattern) if recursive else p.glob(pattern))
    if not matches:
        return f"No files matching '{pattern}' in {p}"

    lines: list[str] = []
    for f in matches[:200]:
        if f.is_file():
            lines.append(f"  {f.relative_to(p)}  ({f.stat().st_size:,} bytes)")
        else:
            lines.append(f"  {f.relative_to(p)}/")
    if len(matches) > 200:
        lines.append(f"  ... and {len(matches) - 200} more")

    return f"{p}\n" + "\n".join(lines)


def move_file(src: str, dst: str) -> str:
    s = Path(src).expanduser()
    d = Path(dst).expanduser()
    if not s.exists():
        return f"Source not found: {s}"
    d.parent.mkdir(parents=True, exist_ok=True)
    s.rename(d)
    return f"Moved {s} → {d}"


def delete_file(path: str, recursive: bool = False) -> str:
    p = Path(path).expanduser()
    if not p.exists():
        return f"Not found: {p}"
    if p.is_dir():
        if recursive:
            shutil.rmtree(p)
            return f"Deleted directory {p} and all its contents."
        return f"{p} is a directory — set recursive=true to delete it and all contents."
    p.unlink()
    return f"Deleted {p}"


def make_slides(
    title: str,
    slides: list[dict],
    output_path: str,
    format: str = "pptx",
) -> str:
    out = Path(output_path).expanduser()
    out.parent.mkdir(parents=True, exist_ok=True)

    if format == "pdf":
        try:
            from fpdf import FPDF
        except ImportError:
            return "ERROR: fpdf2 is not installed. Run: pip install fpdf2"

        pdf = FPDF()
        pdf.set_auto_page_break(auto=False)

        # Title page
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 32)
        pdf.set_y(80)
        pdf.multi_cell(0, 15, title, align="C")

        for s in slides:
            pdf.add_page()
            pdf.set_font("Helvetica", "B", 22)
            pdf.set_y(20)
            pdf.multi_cell(0, 12, s.get("title", ""), align="L")
            pdf.ln(5)
            pdf.set_font("Helvetica", size=14)
            for line in s.get("content", "").split("\n"):
                line = line.strip()
                if line:
                    pdf.multi_cell(0, 8, f"  •  {line}")

        pdf.output(str(out))

    else:
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            return "ERROR: python-pptx is not installed. Run: pip install python-pptx"

        prs = Presentation()

        # Title slide
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = title
        if len(sl.placeholders) > 1:
            sl.placeholders[1].text = ""

        # Content slides
        for s in slides:
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = s.get("title", "")
            tf = sl.placeholders[1].text_frame
            tf.clear()
            lines = [ln.strip() for ln in s.get("content", "").split("\n") if ln.strip()]
            for i, line in enumerate(lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0

        prs.save(str(out))

    kb = out.stat().st_size // 1024
    return f"Saved {format.upper()} ({len(slides)} content slides, {kb} KB) → {out}"


async def install_package(package: str) -> str:
    proc = await asyncio.create_subprocess_exec(
        sys.executable,
        "-m", "pip", "install", package, "--quiet",
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )
    try:
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=120)
    except asyncio.TimeoutError:
        proc.kill()
        return "Timeout: pip install took more than 120 seconds."

    out = stdout.decode(errors="replace")
    err = stderr.decode(errors="replace")
    if proc.returncode == 0:
        return f"Successfully installed {package}.\n{out}".strip()
    return f"Failed to install {package} (exit {proc.returncode}).\n{err}".strip()


async def rag_search(query: str, *, config, store, ollama) -> str:
    from offline_rag.retrieval import retrieve

    result = await retrieve(query, config, store, ollama)
    if not result.context_text:
        return "No relevant documents found in your local index."
    sources = ", ".join(result.sources[:5])
    return f"Sources: {sources}\n\n{result.context_text}"


async def index_documents(paths: list[str], *, config, store, ollama) -> str:
    from offline_rag.indexer import FileIndexer

    indexer = FileIndexer(config, store, ollama)
    counts = await indexer.index_paths(paths)
    return (
        f"Indexed {counts.get('indexed', 0)} files, "
        f"skipped {counts.get('skipped', 0)} unchanged, "
        f"{counts.get('errors', 0)} errors."
    )


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

async def execute_tool(
    name: str,
    args: dict,
    *,
    config=None,
    store=None,
    ollama=None,
) -> str:
    match name:
        case "web_search":
            return await web_search(**args)
        case "web_fetch":
            return await web_fetch(**args)
        case "run_python":
            return await run_python(**args)
        case "run_shell":
            return await run_shell(**args)
        case "read_file":
            return read_file(**args)
        case "write_file":
            return write_file(**args)
        case "list_files":
            return list_files(**args)
        case "move_file":
            return move_file(**args)
        case "delete_file":
            return delete_file(**args)
        case "make_slides":
            return make_slides(**args)
        case "rag_search":
            return await rag_search(config=config, store=store, ollama=ollama, **args)
        case "install_package":
            return await install_package(**args)
        case "index_documents":
            return await index_documents(config=config, store=store, ollama=ollama, **args)
        case _:
            return f"Unknown tool: {name}"
